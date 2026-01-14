"""
pdf2ppt_qwenvl workflow
~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
基于 slides PDF，融合 VLM (Qwen-VL-OCR) 替代传统 PaddleOCR：
1. 将 PDF 每页渲染为 PNG
2. 对每页图片用 VLM (ImageTextBBoxAgent) 做文字识别与定位 (替代 PaddleOCR)
3. 对每页图片用 MinerU 做版面分析（区分 Text vs Image/Table）
4. 对每页图片用 SAM 做图标 / 图块分割
5. AI 背景编辑：
   - 基于 VLM 提取的 bbox 生成 mask (或利用 VLM 调试阶段的 no_text 图)
   - 调用 Inpainting API：填补 mask 之后的白色区域，结合背景颜色做 Inpainting
6. 智能合并与 PPT 生成：
   - 结合 MinerU (版面), SAM (图标), VLM (文字) 结果生成 PPT。
"""

from __future__ import annotations
import os
import asyncio
from pathlib import Path
from typing import List, Dict, Any, Optional
from collections import Counter
import copy

import cv2
import numpy as np
import fitz  # PyMuPDF
import yaml
from PIL import Image

from dataflow_agent.workflow.registry import register
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.logger import get_logger

from dataflow_agent.state import Paper2FigureState 
from dataflow_agent.utils import get_project_root
from dataflow_agent.agentroles import create_vlm_agent

# Tools
from dataflow_agent.toolkits.imtool.sam_tool import segment_layout_boxes, segment_layout_boxes_server, free_sam_model
from dataflow_agent.toolkits.imtool.bg_tool import local_tool_for_bg_remove, free_bg_rm_model
from dataflow_agent.toolkits.imtool.mineru_tool import recursive_mineru_layout
from dataflow_agent.toolkits.imtool.req_img import generate_or_edit_and_save_image_async
from dataflow_agent.toolkits.imtool import ppt_tool

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

log = get_logger(__name__)

# Load configuration from yaml
def load_server_config():
    root = get_project_root()
    config_path = root / "conf" / "model_servers.yaml"
    if not config_path.exists():
        log.warning(f"Config file not found at {config_path}, using defaults.")
        return {}
    try:
        with open(config_path, "r") as f:
            return yaml.safe_load(f) or {}
    except Exception as e:
        log.error(f"Failed to load config: {e}")
        return {}

SERVER_CONFIG = load_server_config()

def get_closest_aspect_ratio(w: int, h: int) -> str:
    """
    计算最接近的合法 Gemini 比例
    """
    valid_ratios = ['1:1', '2:3', '3:2', '3:4', '4:3', '4:5', '5:4', '9:16', '16:9', '21:9']
    target_ratio = w / h
    
    best_ratio = '16:9' # default
    min_diff = float('inf')
    
    for r_str in valid_ratios:
        rw, rh = map(int, r_str.split(':'))
        curr_ratio = rw / rh
        diff = abs(target_ratio - curr_ratio)
        if diff < min_diff:
            min_diff = diff
            best_ratio = r_str
            
    return best_ratio

# Helper to construct URLs
def get_sam_urls():
    if os.environ.get("SAM_SERVER_URLS"):
        return os.environ.get("SAM_SERVER_URLS").split(",")
    sam_cfg = SERVER_CONFIG.get("sam", {})
    instances = sam_cfg.get("instances", [])
    if instances:
        urls = []
        for inst in instances:
            for port in inst.get("ports", []):
                urls.append(f"http://127.0.0.1:{port}")
        if urls:
            return urls
    return ["http://localhost:8021", "http://localhost:8022","http://localhost:8023"]

SAM_SERVER_URLS = get_sam_urls()

def _ensure_result_path(state: Paper2FigureState) -> str:
    raw = getattr(state, "result_path", None)
    if raw:
        return raw
    root = get_project_root()
    ts = int(__import__("time").time())
    base_dir = (root / "outputs" / "pdf2ppt_qwenvl" / str(ts)).resolve()
    base_dir.mkdir(parents=True, exist_ok=True)
    state.result_path = str(base_dir)
    return state.result_path

def _run_sam_on_pages(image_paths: List[str], base_dir: str) -> List[Dict[str, Any]]:
    results: List[Dict[str, Any]] = []
    sam_ckpt = f"{get_project_root()}/sam_b.pt"

    for page_idx, img_path in enumerate(image_paths):
        img_path_obj = Path(img_path)
        if not img_path_obj.exists():
            results.append({"page_idx": page_idx, "layout_items": []})
            continue

        out_dir = Path(base_dir) / "layout_items" / f"page_{page_idx+1:03d}"
        out_dir.mkdir(parents=True, exist_ok=True)

        try:
            layout_items = segment_layout_boxes_server(
                image_path=str(img_path_obj),
                output_dir=str(out_dir),
                server_urls=SAM_SERVER_URLS,
                checkpoint=sam_ckpt,
                min_area=200,
                min_score=0.0,
                iou_threshold=0.4,
                top_k=25,
                nms_by="mask",
            )
        except Exception as e:
            log.error(f"[pdf2ppt_qwenvl] Remote SAM failed: {e}. Fallback to local.")
            layout_items = segment_layout_boxes(
                image_path=str(img_path_obj),
                output_dir=str(out_dir),
                checkpoint=sam_ckpt,
                min_area=200,
                min_score=0.0,
                iou_threshold=0.4,
                top_k=25,
                nms_by="mask",
            )
            
        try:
            pil_img = Image.open(str(img_path_obj))
            w, h = pil_img.size
        except Exception:
            w, h = 1024, 768

        for it in layout_items:
            bbox = it.get("bbox")
            if bbox and len(bbox) == 4:
                x1n, y1n, x2n, y2n = bbox
                x1 = int(round(x1n * w))
                y1 = int(round(y1n * h))
                x2 = int(round(x2n * w))
                y2 = int(round(y2n * h))
                if x2 > x1 and y2 > y1:
                    it["bbox_px"] = [x1, y1, x2, y2]

        results.append({"page_idx": page_idx, "layout_items": layout_items})

    try:
        free_sam_model(checkpoint=sam_ckpt)
    except Exception:
        pass

    return results

@register("pdf2ppt_qwenvl")
def create_pdf2ppt_qwenvl_graph() -> GenericGraphBuilder:
    """
    Workflow factory: dfa run --wf pdf2ppt_qwenvl
    """
    builder = GenericGraphBuilder(state_model=Paper2FigureState, entry_point="_start_")

    def _init_result_path(state: Paper2FigureState) -> Paper2FigureState:
        _ensure_result_path(state)
        return state

    async def pdf_to_images_node(state: Paper2FigureState) -> Paper2FigureState:
        if state.request.input_type == "FIGURE":
            img_path = state.request.input_content
            log.info(f"[pdf2ppt_qwenvl] FIGURE mode: using input image {img_path}")
            state.use_ai_edit = True
            if img_path and os.path.exists(img_path):
                state.slide_images = [img_path]
            else:
                log.error(f"[pdf2ppt_qwenvl] FIGURE mode: image not found {img_path}")
            return state

        pdf_path = getattr(state, "pdf_file", None)
        if not pdf_path:
            log.error("[pdf2ppt_qwenvl] state.pdf_file is empty")
            return state

        base_dir = Path(_ensure_result_path(state))
        img_dir = base_dir / "slides_png"
        image_paths = ppt_tool.pdf_to_images(pdf_path, str(img_dir))
        state.slide_images = image_paths
        return state

    # --- 新增：VLM 节点 (替代原 OCR) ---
    async def vlm_recognition_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        使用 VLM (ImageTextBBoxAgent) 提取文本和 bbox。
        结果写入 state.vlm_pages。
        """
        image_paths: List[str] = getattr(state, "slide_images", []) or []
        if not image_paths:
            log.warning("[pdf2ppt_qwenvl][vlm] no slide_images")
            state.vlm_pages = []
            return state

        async def _process_single_image(page_idx: int, img_path: str) -> Dict[str, Any]:
            try:
                # 显式传递 result_path，确保 agent 内部能访问
                temp_state = copy.copy(state)
                temp_state.result_path = state.result_path
                
                # Retry loop for VLM execution
                max_retries = 3
                bbox_res = []
                
                for attempt in range(max_retries):
                    try:
                        agent = create_vlm_agent(
                            name="ImageTextBBoxAgent",
                            model_name=getattr(state.request, "vlm_model", "qwen-vl-ocr-2025-11-20"),
                            chat_api_url=getattr(state.request, "chat_api_url", None),
                            vlm_mode="understanding",
                            additional_params={
                                "input_image": img_path
                            }
                        )

                        log.info(f"[pdf2ppt_qwenvl][VLM] page#{page_idx+1} attempt {attempt+1}/{max_retries}...")
                        new_state = await agent.execute(temp_state)
                        bbox_res = getattr(new_state, "bbox_result", [])
                        
                        # Basic validation: ensure we got a list and it's not empty (unless image is truly blank)
                        # Here we assume a successful parse returns a list. If it failed to parse, base_agent usually returns error dict or empty.
                        # We can check if new_state has error info or if bbox_res is valid.
                        if isinstance(bbox_res, list):
                            # Success
                            break
                        else:
                            log.warning(f"[pdf2ppt_qwenvl][VLM] page#{page_idx+1} attempt {attempt+1} got invalid result: {type(bbox_res)}")
                            
                    except Exception as e:
                        log.warning(f"[pdf2ppt_qwenvl][VLM] page#{page_idx+1} attempt {attempt+1} failed: {e}")
                        if attempt == max_retries - 1:
                            raise e
                        # Continue to retry
                        await asyncio.sleep(1)

                if not isinstance(bbox_res, list):
                    bbox_res = []
                
                # 修正 bbox 归一化 (0-1000 -> 0-1)
                # 并生成 "no_text" mask 图，供后续 Inpainting 使用
                processed_items = []
                
                # 读取原图尺寸
                pil_img = Image.open(img_path)
                w, h = pil_img.size
                
                # 准备生成 no_text 图
                mask_img = cv2.cvtColor(np.array(pil_img), cv2.COLOR_RGB2BGR)
                
                VLM_SCALE = 1000.0
                
                for it in bbox_res:
                    # 处理 rotate_rect
                    if "rotate_rect" in it and "bbox" not in it:
                        try:
                            rr = it["rotate_rect"]
                            if isinstance(rr, list) and len(rr) == 5:
                                cx, cy, rw, rh, angle = rr
                                rect = ((float(cx), float(cy)), (float(rw), float(rh)), float(angle))
                                box = cv2.boxPoints(rect)
                                x_min = np.min(box[:, 0])
                                x_max = np.max(box[:, 0])
                                y_min = np.min(box[:, 1])
                                y_max = np.max(box[:, 1])
                                
                                it["bbox"] = [
                                    max(0.0, min(1.0, y_min / VLM_SCALE)),
                                    max(0.0, min(1.0, x_min / VLM_SCALE)),
                                    max(0.0, min(1.0, y_max / VLM_SCALE)),
                                    max(0.0, min(1.0, x_max / VLM_SCALE))
                                ]
                        except Exception:
                            pass
                    
                    if "bbox" in it:
                        processed_items.append(it)
                        # 在 mask_img 上将文字区域涂白
                        y1_n, x1_n, y2_n, x2_n = it["bbox"]
                        x1 = int(x1_n * w)
                        y1 = int(y1_n * h)
                        x2 = int(x2_n * w)
                        y2 = int(y2_n * h)
                        # 稍微扩大一点 mask 区域以覆盖完全
                        pad = 2
                        x1 = max(0, x1 - pad)
                        y1 = max(0, y1 - pad)
                        x2 = min(w, x2 + pad)
                        y2 = min(h, y2 + pad)
                        
                        cv2.rectangle(mask_img, (x1, y1), (x2, y2), (255, 255, 255), -1)

                # 保存 no_text 图片
                base_dir = Path(_ensure_result_path(state))
                debug_dir = base_dir / "vlm_debug"
                debug_dir.mkdir(parents=True, exist_ok=True)
                no_text_path = debug_dir / f"page_{page_idx+1:03d}_no_text.png"
                cv2.imwrite(str(no_text_path), mask_img)
                
                log.info(f"[pdf2ppt_qwenvl][VLM] page#{page_idx+1} items={len(processed_items)}, saved mask to {no_text_path}")
                
                return {
                    "page_idx": page_idx,
                    "path": img_path,
                    "vlm_data": processed_items,
                    "no_text_path": str(no_text_path)
                }

            except Exception as e:
                log.error(f"[pdf2ppt_qwenvl][VLM] page#{page_idx+1} failed: {e}")
                return {
                    "page_idx": page_idx,
                    "path": img_path,
                    "vlm_data": [],
                    "error": str(e)
                }

        tasks = [_process_single_image(i, p) for i, p in enumerate(image_paths)]
        results = await asyncio.gather(*tasks)
        state.vlm_pages = results
        return state

    async def slides_mineru_node(state: Paper2FigureState) -> Paper2FigureState:
        """MinerU 版面分析 (并行优化)"""
        image_paths: List[str] = getattr(state, "slide_images", []) or []
        if not image_paths:
            return state

        base_dir = Path(_ensure_result_path(state))
        mineru_dir = base_dir / "mineru_pages"
        mineru_dir.mkdir(parents=True, exist_ok=True)
        port = getattr(getattr(state, "request", None), "mineru_port", 8010)

        async def _process_mineru_page(page_idx: int, img_path: str) -> Dict[str, Any]:
            try:
                out_dir = mineru_dir / f"page_{page_idx+1:03d}"
                out_dir.mkdir(parents=True, exist_ok=True)

                mineru_items = await recursive_mineru_layout(
                    image_path=str(img_path),
                    port=port,
                    max_depth=3,
                    output_dir=str(out_dir),
                )
                return {
                    "page_idx": page_idx,
                    "blocks": mineru_items,
                    "path": img_path,
                    "mineru_output_dir": str(out_dir)
                }
            except Exception as e:
                log.error(f"[pdf2ppt_qwenvl][MinerU] page#{page_idx+1} failed: {e}")
                return {
                    "page_idx": page_idx,
                    "blocks": [],
                    "path": img_path
                }

        tasks = [_process_mineru_page(i, p) for i, p in enumerate(image_paths)]
        results = await asyncio.gather(*tasks)
        
        # 按 page_idx 排序确保顺序
        state.mineru_pages = sorted(results, key=lambda x: x["page_idx"])
        return state

    async def slides_sam_node(state: Paper2FigureState) -> Paper2FigureState:
        """SAM 图标分割"""
        image_paths: List[str] = getattr(state, "slide_images", []) or []
        if not image_paths:
            return state
        base_dir = _ensure_result_path(state)
        sam_pages = await asyncio.to_thread(_run_sam_on_pages, image_paths, base_dir)
        state.sam_pages = sam_pages
        return state

    async def slides_layout_bg_remove_node(state: Paper2FigureState, sam_pages: List[Dict[str, Any]] = None) -> Paper2FigureState:
        """SAM 结果背景移除"""
        if sam_pages is None:
            sam_pages = getattr(state, "sam_pages", []) or []
        if not sam_pages:
            return state

        base_dir = Path(_ensure_result_path(state))
        icons_dir = base_dir / "sam_icons"
        icons_dir.mkdir(parents=True, exist_ok=True)
        model_path = getattr(getattr(state, "request", None), "bg_rm_model", None)

        def _sync_bg_remove():
            processed = 0
            for p in sam_pages:
                page_idx = p.get("page_idx", 0)
                for it in p.get("layout_items", []):
                    png_path = it.get("png_path")
                    if not png_path or not os.path.exists(png_path): continue
                    
                    try:
                        original_stem = Path(png_path).stem
                        output_filename = f"page_{page_idx+1:03d}_{original_stem}_bg_removed.png"
                        output_path = icons_dir / output_filename
                        
                        req = {"image_path": png_path, "output_dir": str(icons_dir)}
                        if model_path: req["model_path"] = model_path
                        
                        fg_path = local_tool_for_bg_remove(req)
                        
                        if fg_path and os.path.exists(fg_path):
                            fg_path_obj = Path(fg_path)
                            if fg_path_obj.name != output_filename:
                                new_fg_path = fg_path_obj.parent / output_filename
                                fg_path_obj.rename(new_fg_path)
                                fg_path = str(new_fg_path)
                            it["fg_png_path"] = fg_path
                        else:
                            it["fg_png_path"] = png_path
                        processed += 1
                    except Exception:
                        it["fg_png_path"] = png_path
            
            try:
                if model_path: free_bg_rm_model(model_path=model_path)
            except Exception: pass
            return processed

        await asyncio.to_thread(_sync_bg_remove)
        state.sam_pages = sam_pages
        return state

    # --- 并行处理节点 ---
    async def parallel_processing_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        并行执行：
        1. VLM (OCR)
        2. MinerU
        3. SAM + BgRemove
        """
        import copy
        import time
        start_time = time.time()
        
        async def vlm_branch():
            branch_state = copy.copy(state)
            result = await vlm_recognition_node(branch_state)
            return ("vlm", result)
        
        async def mineru_branch():
            branch_state = copy.copy(state)
            result = await slides_mineru_node(branch_state)
            return ("mineru", result)
        
        async def sam_branch():
            branch_state = copy.copy(state)
            branch_state = await slides_sam_node(branch_state)
            sam_pages = getattr(branch_state, "sam_pages", [])
            branch_state = await slides_layout_bg_remove_node(branch_state, sam_pages=sam_pages)
            return ("sam", branch_state)
        
        results = await asyncio.gather(
            vlm_branch(),
            mineru_branch(),
            sam_branch(),
            return_exceptions=True
        )
        
        for r in results:
            if isinstance(r, Exception):
                log.error(f"[pdf2ppt_qwenvl] Branch failed: {r}")
                continue
            branch_name, branch_state = r
            if branch_name == "vlm":
                state.vlm_pages = getattr(branch_state, "vlm_pages", [])
            elif branch_name == "mineru":
                state.mineru_pages = getattr(branch_state, "mineru_pages", [])
            elif branch_name == "sam":
                state.sam_pages = getattr(branch_state, "sam_pages", [])
        
        log.info(f"[pdf2ppt_qwenvl] Parallel processing finished in {time.time() - start_time:.2f}s")
        return state

    async def slides_ppt_generation_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        生成 PPT：
        1. 整合 VLM, MinerU, SAM 结果
        2. 调用 AI Inpainting (利用 no_text mask 图)
        3. 渲染页面
        """
        vlm_pages = getattr(state, "vlm_pages", []) or []
        sam_pages = getattr(state, "sam_pages", []) or []
        mineru_pages = getattr(state, "mineru_pages", []) or []

        if not vlm_pages:
            log.error("[pdf2ppt_qwenvl] no vlm_pages, abort PPT generation")
            return state

        # Indexing
        sam_dict = {p.get("page_idx", 0): p.get("layout_items", []) for p in sam_pages}
        mineru_dict = {}
        for p in mineru_pages:
            mineru_dict[p.get("page_idx", 0)] = p

        prs = Presentation()
        prs.slide_width = Inches(ppt_tool.SLIDE_W_IN)
        prs.slide_height = Inches(ppt_tool.SLIDE_H_IN)
        slide_w_emu = prs.slide_width
        slide_h_emu = prs.slide_height

        base_dir = Path(_ensure_result_path(state))

        # API helper
        async def _call_image_api_with_retry(coro_factory, retries=3):
            for i in range(retries):
                try:
                    await coro_factory()
                    return True
                except Exception as e:
                    if i == retries - 1: log.error(f"Image API failed: {e}")
                    await asyncio.sleep(1)
            return False

        # --- 准备渲染数据 & AI Inpainting ---
        pages_render_data = []
        ai_coroutines = []
        
        # 限制并发数为 1 (串行处理) 以避免 413 错误及并发限制
        sem = asyncio.Semaphore(3)
        
        req_cfg = getattr(state, "request", None) or {}
        if not isinstance(req_cfg, dict): req_cfg = req_cfg.__dict__ if hasattr(req_cfg, "__dict__") else {}
        api_key = req_cfg.get("api_key") or os.getenv("DF_API_KEY")
        api_url = req_cfg.get("chat_api_url") or "https://api.apiyi.com"
        model_name = req_cfg.get("gen_fig_model") or "gemini-3-pro-image-preview"

        # 辅助几何函数
        def _bbox_area(bbox): return max(0, bbox[2]-bbox[0]) * max(0, bbox[3]-bbox[1])
        def _get_intersection_area(b1, b2):
            x1,y1,x2,y2 = max(b1[0],b2[0]), max(b1[1],b2[1]), min(b1[2],b2[2]), min(b1[3],b2[3])
            return max(0, x2-x1) * max(0, y2-y1)
        def _is_inside(inner, outer, th=0.9):
            ia = _bbox_area(inner)
            return (ia > 0) and ((_get_intersection_area(inner, outer) / ia) >= th)
        
        for pinfo in vlm_pages:
            page_idx = pinfo.get("page_idx", 0)
            img_path = pinfo.get("path")
            no_text_path = pinfo.get("no_text_path")
            vlm_data = pinfo.get("vlm_data", [])

            if not img_path or not os.path.exists(img_path): continue

            try:
                pil_img = Image.open(img_path)
                w0, h0 = pil_img.size
            except Exception: continue

            # 1. MinerU Image Zones
            mineru_data = mineru_dict.get(page_idx, {})
            mineru_blocks = mineru_data.get("blocks", [])
            image_zones = []
            
            # 简化版图片查找逻辑（复用 parallel 版的核心思想）
            sub_images_dir = None
            if mineru_data.get("mineru_output_dir"):
                try:
                    # 简单尝试找 sub_images
                    possibles = list(Path(mineru_data["mineru_output_dir"]).rglob("sub_images"))
                    if possibles: sub_images_dir = possibles[0]
                except Exception: pass

            for idx, blk in enumerate(mineru_blocks):
                btype = (blk.get("type") or "").lower()
                bbox = blk.get("bbox") # norm
                if not bbox or len(bbox)!=4: continue
                if btype in ['image', 'figure', 'table', 'formula']:
                    x1, y1, x2, y2 = int(bbox[0]*w0), int(bbox[1]*h0), int(bbox[2]*w0), int(bbox[3]*h0)
                    px_bbox = [x1, y1, x2, y2]
                    
                    img_path_found = None
                    # 尝试从 blk['img_path'] 或 sub_images 或 裁剪
                    if blk.get("img_path") and os.path.exists(blk["img_path"]):
                        img_path_found = blk["img_path"]
                    elif sub_images_dir:
                        # 简化匹配：尝试找对应索引
                        prefix = f"_{idx}_" # 假设文件名含索引，这比较脆弱，降级为裁剪
                        pass 
                    
                    if not img_path_found:
                        fb_dir = base_dir / "mineru_fallback" / f"p{page_idx}"
                        fb_dir.mkdir(parents=True, exist_ok=True)
                        save_p = fb_dir / f"blk_{idx}.png"
                        if not save_p.exists():
                            try: pil_img.crop((x1,y1,x2,y2)).save(save_p)
                            except: pass
                        img_path_found = str(save_p)
                    
                    if img_path_found:
                        image_zones.append({"bbox": px_bbox, "type": btype, "img_path": img_path_found})

            # 2. VLM Text Filtering (Filter text inside images)
            final_text_lines = []
            for it in vlm_data:
                # it: {'bbox': [y1n, x1n, y2n, x2n], 'text': ...}  (0-1 norm)
                bbox_n = it.get("bbox")
                if not bbox_n: continue
                y1n, x1n, y2n, x2n = bbox_n
                x1, y1, x2, y2 = int(x1n*w0), int(y1n*h0), int(x2n*w0), int(y2n*h0)
                l_bbox = [x1, y1, x2, y2]
                
                is_in_image = False
                # [MODIFIED] Keep all VLM text, even if inside image zones
                # for z in image_zones:
                #     if _is_inside(l_bbox, z["bbox"]): 
                #         is_in_image = True
                #         break
                
                if not is_in_image:
                    # 估算字号
                    raw_pt_obj = ppt_tool.estimate_font_pt(l_bbox, img_h_px=h0, body_h_px=None)
                    raw_pt = raw_pt_obj.pt if hasattr(raw_pt_obj, "pt") else raw_pt_obj
                    
                    # 简单判断 title (基于字号或位置，这里简化)
                    l_type = "body"
                    if raw_pt > 18: l_type = "title" # 简单阈值，可改进
                    
                    final_text_lines.append((l_bbox, it.get("text",""), 1.0, l_type, raw_pt))

            # 3. SAM Icons Filtering
            raw_sam = sam_dict.get(page_idx, [])
            final_sam = []
            for item in raw_sam:
                s_bbox = item.get("bbox_px")
                if not s_bbox: continue
                # Filter if inside Image Zone
                if any(_is_inside(s_bbox, z["bbox"], 0.6) for z in image_zones): continue
                # Filter if overlaps with Text
                if any(_is_inside(line[0], s_bbox) for line in final_text_lines): continue
                
                final_sam.append(item)

            # 4. AI Inpainting
            clean_bg_path = base_dir / "clean_bg" / f"bg_{page_idx+1:03d}.png"
            clean_bg_path.parent.mkdir(parents=True, exist_ok=True)
            
            # 使用用户指定的 Prompt
            inpainting_prompt = "使用背景颜色，填充图里被mask的白色部分，去掉全部文字！"
            
            if state.use_ai_edit and api_key and no_text_path and os.path.exists(no_text_path):
                # 动态计算比例
                ratio_str = "16:9"
                try:
                    with Image.open(no_text_path) as tmp_img:
                        ratio_str = get_closest_aspect_ratio(tmp_img.width, tmp_img.height)
                except Exception:
                    pass

                async def _run_ai(
                    prompt=inpainting_prompt, 
                    save_path=str(clean_bg_path), 
                    img_path=no_text_path, 
                    ratio=ratio_str
                ):
                    async with sem:
                        await _call_image_api_with_retry(
                            lambda: generate_or_edit_and_save_image_async(
                                prompt=prompt,
                                save_path=save_path,
                                api_url=api_url,
                                api_key=api_key,
                                model=model_name,
                                use_edit=True,
                                image_path=img_path,  # 传入 no_text 图作为原图进行编辑
                                aspect_ratio=ratio,   # 动态计算的比例
                                resolution="2K"       # 默认分辨率
                            )
                        )
                ai_coroutines.append(_run_ai())
            else:
                # 降级：直接用 no_text_path (虽然有白块) 或者 原图
                if no_text_path and os.path.exists(no_text_path):
                     import shutil
                     try: shutil.copy(no_text_path, clean_bg_path)
                     except: pass

            pages_render_data.append({
                "page_idx": page_idx,
                "scale_x": slide_w_emu / w0,
                "scale_y": slide_h_emu / h0,
                "clean_bg_path": str(clean_bg_path),
                "image_zones": image_zones,
                "final_sam": final_sam,
                "final_text": final_text_lines
            })

        # Execute AI
        if ai_coroutines:
            log.info(f"[pdf2ppt_qwenvl] Executing {len(ai_coroutines)} Inpainting tasks...")
            await asyncio.gather(*ai_coroutines, return_exceptions=True)

        # 渲染 PPT
        for p_data in pages_render_data:
            scale_x = p_data["scale_x"]
            scale_y = p_data["scale_y"]
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Background
            if os.path.exists(p_data["clean_bg_path"]):
                try: slide.shapes.add_picture(p_data["clean_bg_path"], 0, 0, prs.slide_width, prs.slide_height)
                except: pass
            
            # MinerU Images
            for z in p_data["image_zones"]:
                if os.path.exists(z["img_path"]):
                    bx = z["bbox"]
                    slide.shapes.add_picture(z["img_path"], 
                        ppt_tool.px_to_emu(bx[0], scale_x), ppt_tool.px_to_emu(bx[1], scale_y),
                        ppt_tool.px_to_emu(bx[2]-bx[0], scale_x), ppt_tool.px_to_emu(bx[3]-bx[1], scale_y))
            
            # SAM Icons
            for s in p_data["final_sam"]:
                path = s.get("fg_png_path") or s.get("png_path")
                if path and os.path.exists(path):
                    bx = s["bbox_px"]
                    slide.shapes.add_picture(path,
                        ppt_tool.px_to_emu(bx[0], scale_x), ppt_tool.px_to_emu(bx[1], scale_y),
                        ppt_tool.px_to_emu(bx[2]-bx[0], scale_x), ppt_tool.px_to_emu(bx[3]-bx[1], scale_y))

            # Text
            for line in p_data["final_text"]:
                bbox, text, _, l_type, raw_pt = line
                left = ppt_tool.px_to_emu(bbox[0], scale_x)
                top = ppt_tool.px_to_emu(bbox[1], scale_y)
                w = ppt_tool.px_to_emu(bbox[2]-bbox[0], scale_x)
                h = ppt_tool.px_to_emu(bbox[3]-bbox[1], scale_y)
                
                tb = slide.shapes.add_textbox(left, top, w, h)
                p = tb.text_frame.paragraphs[0]
                p.text = text
                p.font.size = Pt(raw_pt if raw_pt > 5 else 12)
                p.font.bold = (l_type == "title")
                p.font.color.rgb = RGBColor(0,0,0)

        out_path = base_dir / "pdf2ppt_qwenvl_output.pptx"
        prs.save(str(out_path))
        state.ppt_path = str(out_path)
        log.info(f"[pdf2ppt_qwenvl] PPT Generated: {out_path}")
        return state

    nodes = {
        "_start_": _init_result_path,
        "pdf_to_images": pdf_to_images_node,
        "parallel_processing": parallel_processing_node,
        "slides_ppt_generation": slides_ppt_generation_node,
        "_end_": lambda s: s,
    }

    edges = [
        ("pdf_to_images", "parallel_processing"),
        ("parallel_processing", "slides_ppt_generation"),
        ("slides_ppt_generation", "_end_"),
    ]

    builder.add_nodes(nodes).add_edges(edges)
    builder.add_edge("_start_", "pdf_to_images")
    return builder
