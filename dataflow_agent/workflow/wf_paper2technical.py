"""
paper2technical workflow
~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~~
生成时间: 2025-12-07 23:36:51

1. 在 **TOOLS** 区域定义需要暴露给 Prompt 的前置工具
2. 在 **NODES**  区域实现异步节点函数 (await-able)
3. 在 **EDGES**  区域声明有向边
4. 最后返回 builder.compile() 或 GenericGraphBuilder
"""

from __future__ import annotations
import json
import time
from pathlib import Path
import re

from dataflow_agent.state import Paper2FigureState
from dataflow_agent.graphbuilder.graph_builder import GenericGraphBuilder
from dataflow_agent.workflow.registry import register
from dataflow_agent.agentroles import create_graph_agent, create_react_agent, create_simple_agent
from dataflow_agent.toolkits.tool_manager import get_tool_manager
from dataflow_agent.toolkits.imtool.bg_tool import (
    local_tool_for_svg_render,
    local_tool_for_raster_to_svg,
)
from dataflow_agent.toolkits.imtool.mineru_tool import svg_to_emf
from dataflow_agent.utils import get_project_root
from dataflow_agent.logger import get_logger
log = get_logger(__name__)


def _ensure_result_path(state: Paper2FigureState) -> str:
    """
    统一本次 workflow 的根输出目录：
    - 如果 state.result_path 已存在（通常由调用方传入，形如 时间戳+编码），直接使用；
    - 否则：使用 get_project_root() / "outputs" / "paper2tec" / <timestamp>，
      并回写到 state.result_path，确保后续节点共享同一目录，避免数据串台。
    """
    raw = getattr(state, "result_path", None)
    if raw:
        return raw

    root = get_project_root()
    ts = int(time.time())
    base_dir = (root / "outputs" / "paper2tec" / str(ts)).resolve()
    base_dir.mkdir(parents=True, exist_ok=True)
    state.result_path = str(base_dir)
    return state.result_path


@register("paper2technical")
def create_paper2technical_graph() -> GenericGraphBuilder:  # noqa: N802
    """
    Workflow factory: dfa run --wf paper2technical
    """
    # 使用 Paper2FigureState，复用其中的 paper_file / paper_idea / fig_desc 等字段，
    # 这里不做图像生成和抠图，只负责“技术路线图”的 SVG + PPT 逻辑。
    builder = GenericGraphBuilder(
        state_model=Paper2FigureState,
        entry_point="_start_",        # 入口统一为 _start_，再由路由函数分发
    )

    # ----------------------------------------------------------------------
    # TOOLS (pre_tool definitions)
    # ----------------------------------------------------------------------
    # 1) 提供给 paper_idea_extractor 的 PDF 内容（标题 + 前几页正文）
    @builder.pre_tool("paper_content", "paper_idea_extractor")
    def _get_paper_content(state: Paper2FigureState):
        """
        前置工具: 读取论文 PDF 的标题和前若干页内容，供 paper_idea_extractor 节点使用。

        - 作用: 为大模型提供足够的上下文，让其抽取论文中的技术路线/实验流程关键信息。
        - 输出: 一个字符串，包含论文标题 + 前若干页文本。
        """
        import fitz  # PyMuPDF
        import PyPDF2

        pdf_path = state.paper_file
        if not pdf_path:
            log.warning("paper_file 为空，无法读取 PDF 内容")
            return ""

        try:
            with open(pdf_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                paper_title = reader.metadata.get("/Title", "Unknown Title")
        except Exception:
            paper_title = "Unknown Title"

        try:
            doc = fitz.open(pdf_path)
        except Exception as e:
            log.error(f"打开 PDF 失败: {e}")
            return f"The title of the paper is {paper_title}"

        text_parts: list[str] = []
        # 读取前 10 页内容，通常技术路线、整体框架会在前几页出现
        for page_idx in range(min(10, len(doc))):
            page = doc.load_page(page_idx)
            text_parts.append(page.get_text("text") or "")

        content = "\n".join(text_parts).strip()
        final_text = (
            f"The title of the paper is {paper_title}\n\n"
            f"Here are the first 10 pages of the paper:\n{content}"
        )
        log.info("paper_content 提取完成")
        return final_text

    @builder.pre_tool("paper_idea", "technical_route_desc_generator")
    def _get_paper_idea(state: Paper2FigureState):
        """
        前置工具: 为 technical_route_desc_generator 节点暴露论文的核心想法摘要。

        - 在 PDF 模式下，该摘要由 paper_idea_extractor 节点写入 state.paper_idea。
        - 在 TEXT 模式下，可以直接由调用方事先把概要写入 state.paper_idea。
        """
        return state.paper_idea or ""

    @builder.pre_tool("style", "technical_route_desc_generator")
    def _get_paper_idea(state: Paper2FigureState):

        return state.request.style or ""

    # ----------------------------------------------------------------------

    # ==============================================================
    # NODES
    # ==============================================================
    async def paper_idea_extractor_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        节点 1: 从 PDF 中抽取论文的核心思想 / 技术路线相关信息

        - 只在 input_type == "PDF" 时作为入口节点被调用。
        - 基于 pre_tool("paper_content") 提供的标题 + 前若干页内容，
          调用专门的 agent（例如 paper_idea_extractor）生成摘要。
        - 该摘要用于后续技术路线图描述生成。

        输入:
            state.paper_file : 论文 PDF 路径
        输出:
            state.paper_idea : 论文核心思想 / 技术路线要点摘要
            state.agent_results["paper_idea_extractor"] : agent 原始输出
        """
        agent = create_simple_agent("paper_idea_extractor")
        state = await agent.execute(state=state)
        return state

    def _svg_has_cjk(text: str) -> bool:
        """简单判断 SVG 中是否包含中文字符，用于日志和调试。"""
        return bool(re.search(r"[\u4e00-\u9fff]", text))


    def _inject_chinese_font(svg_code: str) -> str:
        """
        如果 SVG 中没有设定中文友好的 font-family，则注入一段全局样式，
        指定一组 CJK 字体作为优先字体。

        注意：字体名请根据实际安装的字体调整。
        """
        if "font-family" in svg_code:
            return svg_code

        idx = svg_code.find(">")
        if idx == -1:
            return svg_code

        style_block = """
  <style type="text/css">
    text, tspan {
      font-family: "Noto Sans CJK SC", "Microsoft YaHei", "SimHei", "SimSun", "WenQuanYi Zen Hei", sans-serif;
    }
  </style>
"""
        return svg_code[: idx + 1] + style_block + svg_code[idx + 1 :]


    async def technical_route_desc_generator_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        节点 2: 技术路线图描述生成器

        - 根据论文摘要（PDF 模式）或用户直接提供的文本描述（TEXT 模式），
          生成“技术路线/实验流程”的结构化自然语言描述或 JSON。
        - 典型内容包括: 各阶段实验步骤、模块之间的依赖关系、输入输出数据流等。

        输入:
            - PDF 模式: state.paper_idea 由 paper_idea_extractor 填充
            - TEXT 模式: 可以事先把文本写入 state.paper_idea 或其他字段
        输出:
            - 建议: 在 agent 内把结果存到 state.fig_desc 或 state.agent_results["technical_route_desc_generator"]
        """
        agent = create_react_agent(
            name="technical_route_desc_generator",
            max_retries=4,
            model_name=getattr(state.request, "technical_model", "claude-haiku-4-5-20251001"),
        )
        state = await agent.execute(state=state)

        # --------------------------------------------------------------
        # 将 LLM 生成的 SVG 源码渲染为实际图像文件，并写入统一的 result_path 目录
        # --------------------------------------------------------------
        svg_code = getattr(state, "figure_tec_svg_content", None)
        if svg_code:
            # 日志：是否包含中文
            if _svg_has_cjk(svg_code):
                log.info("technical_route_desc_generator_node: 检测到 SVG 中包含中文字符")

            # 如果 SVG 未显式指定 font-family，则注入一组中文友好的字体
            svg_code = _inject_chinese_font(svg_code)

            # 确保本次 workflow 的根输出目录已确定
            base_dir = Path(_ensure_result_path(state))
            base_dir.mkdir(parents=True, exist_ok=True)

            timestamp = int(time.time())
            # 同时输出 SVG 源码文件和 PNG 位图
            svg_output_path = str((base_dir / f"technical_route_{timestamp}.svg").resolve())
            png_output_path = str((base_dir / f"technical_route_{timestamp}.png").resolve())

            try:
                # 1) 保存 SVG 源码到 .svg 文件
                Path(svg_output_path).write_text(svg_code, encoding="utf-8")
                state.svg_file_path = svg_output_path

                # 2) 渲染 PNG 供 MinerU 使用
                png_path = local_tool_for_svg_render(
                    {
                        "svg_code": svg_code,
                        "output_path": png_output_path,
                    }
                )
                # 将最终图像路径写回 state.svg_img_path
                state.svg_img_path = png_path
                log.critical(f"[state.svg_img_path]: {state.svg_img_path}")
                log.critical(f"[state.svg_file_path]: {state.svg_file_path}")
            except Exception as e:
                # 渲染或写文件失败时仅记录日志，避免打断整体 workflow
                log.error(f"technical_route_desc_generator_node: SVG 落盘/渲染失败: {e}")

        return state

    async def technical_ppt_generator_node(state: Paper2FigureState) -> Paper2FigureState:
        """
        节点 4: 基于技术路线图 PNG 生成 PPT

        - 根据前面步骤生成的 PNG 整图（state.svg_img_path），
          生成用于展示技术路线图的 PPT。
        - 仅使用 SVG -> PNG 路径插入位图，不再尝试 EMF。
        """
        from pptx import Presentation
        from PIL import Image

        # ✅ 临时提高 PIL 图像大小限制，防止 decompression bomb 错误
        original_max_pixels = Image.MAX_IMAGE_PIXELS
        Image.MAX_IMAGE_PIXELS = None  # 或设置为更大的值，如 500_000_000

        try:
            # 输出目录：统一使用本次 workflow 的根输出目录
            run_root = Path(_ensure_result_path(state))
            output_dir = run_root
            output_dir.mkdir(parents=True, exist_ok=True)

            timestamp = int(time.time())
            ppt_path = output_dir / f"technical_route_{timestamp}.pptx"

            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]

            slide_width = prs.slide_width
            slide_height = prs.slide_height

            # ------------------------------------------------------------------
            # 第 1 页：技术路线图（仅使用 PNG）
            # ------------------------------------------------------------------
            slide = prs.slides.add_slide(blank_slide_layout)

            svg_path = getattr(state, "svg_file_path", None)
            png_path = getattr(state, "svg_img_path", None)

            def _insert_picture(pic_path: str) -> bool:
                """通用插图函数：按 80% 宽度缩放并居中。"""
                try:
                    pic = slide.shapes.add_picture(pic_path, 0, 0)
                except Exception as e:
                    log.error(f"technical_ppt_generator_node: 插入图片失败: {e}")
                    return False

                if pic.width and pic.width > 0:
                    scale = (slide_width * 0.8) / pic.width
                else:
                    scale = 1.0

                pic.width = int(pic.width * scale)
                pic.height = int(pic.height * scale)

                pic.left = int((slide_width - pic.width) / 2)
                pic.top = int((slide_height - pic.height) / 2)
                return True

            # 直接使用 PNG（位图）插入 PPT
            inserted = False

            if png_path:
                try:
                    ok = _insert_picture(png_path)
                    if ok:
                        log.info(
                            "technical_ppt_generator_node: 使用 PNG 插入技术路线图成功: %s",
                            png_path,
                        )
                    else:
                        log.error(
                            "technical_ppt_generator_node: PNG 插入失败，第一页可能为空白"
                        )
                except Exception as e:
                    log.error(
                        "technical_ppt_generator_node: PNG 插入失败，第一页将为空白: %s",
                        e,
                    )

            if (not inserted) and (not png_path):
                log.warning(
                    "technical_ppt_generator_node: svg_file_path / svg_img_path 均为空，"
                    "第一页将为空白"
                )

            # ------------------------------------------------------------------
            # 第 2 页：操作提示页（写上“右键转换成形状”）
            # ------------------------------------------------------------------
            slide2 = prs.slides.add_slide(blank_slide_layout)
            left = int(slide_width * 0.1)
            top = int(slide_height * 0.3)
            width = int(slide_width * 0.8)
            height = int(slide_height * 0.4)

            textbox = slide2.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = "右键转换成形状"

            prs.save(str(ppt_path))
            state.ppt_path = str(ppt_path)
            log.info(f"technical_ppt_generator_node: PPT 已生成: {ppt_path}")
        finally:
            # ✅ 恢复原始限制
            Image.MAX_IMAGE_PIXELS = original_max_pixels

        return state

    # ==============================================================
    # 注册 nodes / edges
    # ==============================================================

    def set_entry_node(state: Paper2FigureState) -> str:
        """
        路由函数: 根据输入类型选择技术路线工作流的入口节点。

        - input_type == "PDF"  : 从 PDF 中抽取论文想法，先走 paper_idea_extractor
        - input_type == "TEXT" : 直接使用调用方提供的文本描述，跳过 PDF 抽取，
                                 从 technical_route_desc_generator 开始
        其他值:
        - 认为是不合法输入，直接结束工作流。
        """
        input_type = getattr(state.request, "input_type", "PDF")
        if input_type == "PDF":
            log.critical("paper2technical: 进入 PDF 流程 (paper_idea_extractor)")
            return "paper_idea_extractor"
        elif input_type == "TEXT":
            log.critical("paper2technical: 进入 TEXT 流程 (technical_route_desc_generator)")
            return "technical_route_desc_generator"
        else:
            log.error(f"paper2technical: Invalid input type: {input_type}")
            return "_end_"

    def _init_result_path(state: Paper2FigureState) -> Paper2FigureState:
        """
        _start_ 节点：确保本次 workflow 有一个统一的 result_path 根目录。
        - 若用户已在 state.result_path 传入自定义目录，则直接使用该目录；
        - 若未传入，则初始化为 get_project_root()/outputs/paper2tec/<timestamp>。
        """
        _ensure_result_path(state)
        return state

    nodes = {
        "_start_": _init_result_path,
        "paper_idea_extractor": paper_idea_extractor_node,
        "technical_route_desc_generator": technical_route_desc_generator_node,
        "technical_ppt_generator": technical_ppt_generator_node,
        "_end_": lambda state: state,  # 终止节点
    }

    # ------------------------------------------------------------------
    # EDGES  (从节点 A 指向节点 B)
    # ------------------------------------------------------------------
    edges = [
        # PDF 流程: 先抽想法，再生成技术路线描述
        ("paper_idea_extractor", "technical_route_desc_generator"),
        # PDF/TEXT 后续流程共用: 描述 -> PPT
        ("technical_route_desc_generator", "technical_ppt_generator"),
        ("technical_ppt_generator", "_end_"),
    ]

    builder.add_nodes(nodes).add_edges(edges).add_conditional_edge("_start_", set_entry_node)
    return builder
